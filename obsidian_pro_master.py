#!/usr/bin/env python3
"""
Obsidian Pro Designer Engine — document → preformatted Obsidian notes.

- Multi-format: MarkItDown (office/HTML/etc.), UTF-8 plain text, PDF dual-path
  (MarkItDown + PyMuPDF), raster images (OCR or ![[embed]] fallback).
- Optional OCR: Tesseract via pytesseract for images and low-text PDFs (see config).
- Cross-check: optional length-ratio warning when MarkItDown vs PyMuPDF diverge.
- SHA-256 hashing, atomic writes, optional worker recycling.
- Chapter spindle + index note; designer lint; optional spaCy entities.
- Partial-write rollback; basename de-duplication.
"""

from __future__ import annotations

import argparse
import csv
import html
import json
import logging
import os
import re
import tempfile
import time
import unicodedata
from dataclasses import dataclass, field
from multiprocessing import Pool, cpu_count
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

import yaml

# Optional heavy deps (import lazily in workers if needed)
try:
    from markitdown import MarkItDown
except ImportError:
    MarkItDown = None  # type: ignore

try:
    import spacy
except ImportError:
    spacy = None  # type: ignore

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None  # type: ignore

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None  # type: ignore
    Image = None  # type: ignore

try:
    import docx  # python-docx
except ImportError:
    docx = None  # type: ignore

try:
    import pptx  # python-pptx
except ImportError:
    pptx = None  # type: ignore

try:
    import openpyxl
except ImportError:
    openpyxl = None  # type: ignore


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

# Raster image extensions (OCR or Obsidian embed fallback)
IMAGE_EXTENSIONS = frozenset(
    {".png", ".jpg", ".jpeg", ".webp", ".gif", ".tiff", ".tif", ".bmp"}
)

DEFAULT_CONFIG: dict[str, Any] = {
    "system_settings": {
        "processing_mode": "sequential",
        "core_tiers": {"ultra_safe": 2, "safe": 4, "performance_plus": 8, "aggressive": 12},
        "selected_tier": "safe",
        "hash_algorithm": "sha256",
        "max_tasks_per_worker": 5,
        "atomic_writes": True,
        "log_file": "process_log.json",
        "error_log": "error_log.csv",
        "low_content_threshold_chars": 200,
    },
    "nlp_settings": {
        "engine": "spacy",
        "model": "en_core_web_sm",
        "extract_entities": ["PERSON", "ORG", "GPE"],
        "user_tags": [],
        "extract_concepts_count": 5,
    },
    "designer_settings": {
        "strict_markdown": True,
        "bold_syntax": "asterisks",
        "list_marker": "-",
        "preserve_unicode_filenames": True,
        "chapter_pattern": r"(?i)^#+\s*(Chapter|Section|Part)\s*\d+",
    },
    "pkm_settings": {
        "split_chapters": True,
        "chapter_patterns": [
            r"(?i)^#+\s*(Chapter|Section|Part)\s*\d+\s*:?.*$",
            r"(?i)^Chapter\s+\d+\s*:?.*$",
            r"(?i)^Section\s+\d+\s*:?.*$",
            r"(?i)^(EPILOGUE|PROLOGUE)\b.*$",
        ],
        "ocr_enabled": False,
        "images_ocr": True,
        "ocr_langs": "eng",
        "ocr_max_pdf_pages": 40,
        "ocr_dpi": 150,
        "pdf_pymupdf_text_fallback": True,
        "pdf_cross_check_markitdown_vs_pymupdf": True,
        "ocr_comparison_threshold": 0.35,
        "prefer_longer_text_on_divergence": True,
        "require_markitdown_for_office_formats": False,
        "fail_if_ocr_enabled_but_unavailable": False,
        "callout_keywords": {},
    },
    "yaml_template": {
        "tags": ["processed"],
        "aliases": [],
        "status": "raw",
    },
}


def load_config(path: Path) -> dict[str, Any]:
    if not path.is_file():
        raise FileNotFoundError(f"Config not found: {path}")
    with path.open("r", encoding="utf-8") as f:
        cfg = json.load(f)
    merged = json.loads(json.dumps(DEFAULT_CONFIG))
    _deep_update(merged, cfg)
    _validate_config(merged)
    return merged


def _deep_update(base: dict, over: dict) -> None:
    for k, v in over.items():
        if k in base and isinstance(base[k], dict) and isinstance(v, dict):
            _deep_update(base[k], v)
        else:
            base[k] = v


def _validate_config(cfg: dict[str, Any]) -> None:
    ss = cfg.get("system_settings", {})
    if ss.get("hash_algorithm", "sha256").lower() not in ("sha256", "sha-256"):
        raise ValueError("system_settings.hash_algorithm must be sha256 for this engine.")
    tier = ss.get("selected_tier", "safe")
    tiers = ss.get("core_tiers", {})
    if tier not in tiers:
        raise ValueError(f"selected_tier '{tier}' missing from core_tiers.")
    mtw = int(ss.get("max_tasks_per_worker", 5))
    if mtw < 1:
        raise ValueError("max_tasks_per_worker must be >= 1")


# ---------------------------------------------------------------------------
# Paths / hashing / atomic IO
# ---------------------------------------------------------------------------

def path_for_open(path: Path) -> str:
    """Use Win32 extended path when very long paths are likely."""
    s = str(path.resolve())
    if os.name == "nt" and len(s) > 220 and not s.startswith("\\\\?\\"):
        return "\\\\?\\" + s
    return s


def file_sha256(filepath: Path) -> str:
    import hashlib

    h = hashlib.sha256()
    with open(path_for_open(filepath), "rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def tesseract_is_available() -> bool:
    if pytesseract is None:
        return False
    try:
        _ = pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def doctor_report(cfg: dict[str, Any]) -> tuple[int, str]:
    """
    Return (exit_code, report_text).

    exit_code:
      0 = looks OK for common use
      2 = missing dependency for requested features
    """
    pkm = cfg.get("pkm_settings", {})
    want_ocr = bool(pkm.get("ocr_enabled", False))
    fail_on_ocr = bool(pkm.get("fail_if_ocr_enabled_but_unavailable", False))
    require_office = bool(pkm.get("require_markitdown_for_office_formats", False))

    lines: list[str] = []

    def row(name: str, ok: bool, fix: str = "") -> None:
        lines.append(f"- {name}: {'OK' if ok else 'MISSING'}" + (f" ({fix})" if (fix and not ok) else ""))

    row("MarkItDown", MarkItDown is not None, "pip install markitdown")
    row("PyMuPDF", fitz is not None, "pip install pymupdf")
    row("python-docx", docx is not None, "pip install python-docx")
    row("python-pptx", pptx is not None, "pip install python-pptx")
    row("openpyxl", openpyxl is not None, "pip install openpyxl")
    row("pytesseract", pytesseract is not None, "pip install pytesseract")
    row("Pillow", Image is not None, "pip install Pillow")
    row("tesseract binary", tesseract_is_available(), "install tesseract-ocr on OS and ensure on PATH")

    ok = True

    if MarkItDown is None and fitz is None:
        ok = False
        lines.append("- PDF: neither MarkItDown nor PyMuPDF available; PDFs cannot be processed.")

    if require_office and MarkItDown is None:
        ok = False
        lines.append("- Office formats: require_markitdown_for_office_formats=true but MarkItDown is missing.")

    if want_ocr and not (pytesseract is not None and Image is not None and tesseract_is_available()):
        if fail_on_ocr:
            ok = False
        lines.append("- OCR: enabled but dependencies incomplete; OCR will be skipped with warnings unless fail_if_ocr_enabled_but_unavailable=true.")

    code = 0 if ok else 2
    return code, "Doctor report:\n" + "\n".join(lines) + "\n"


def atomic_write_text(path: Path, text: str, *, encoding: str = "utf-8") -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp = tempfile.mkstemp(
        suffix=".tmp", prefix=path.name + ".", dir=str(path.parent)
    )
    tmp_path = Path(tmp)
    try:
        with os.fdopen(fd, "w", encoding=encoding, newline="\n") as wf:
            wf.write(text)
        os.replace(str(tmp_path), path_for_open(path))
    finally:
        if tmp_path.exists():
            try:
                tmp_path.unlink()
            except OSError:
                pass


# ---------------------------------------------------------------------------
# Text / designer lint
# ---------------------------------------------------------------------------


def slugify_title(text: str, max_len: int = 80, preserve_unicode: bool = True) -> str:
    t = text.strip()
    if not preserve_unicode:
        t = unicodedata.normalize("NFKD", t)
        t = t.encode("ascii", "ignore").decode("ascii")
    # Strip `#` so headings do not pollute filenames / wikilinks on disk.
    t = re.sub(r'[#<>:"/\\|?*\n\r\t]', "", t)
    t = re.sub(r"\s+", " ", t).strip()
    if not t:
        t = "section"
    return t[:max_len]


def apply_designer_linting(text: str, strict: bool, list_marker: str) -> str:
    if not strict:
        return text.rstrip() + "\n"

    out = text.replace("\r\n", "\n").replace("\r", "\n")

    # Bullet normalization: leading * or + -> -
    out = re.sub(r"^(\s*)[*+]\s+", rf"\1{list_marker} ", out, flags=re.MULTILINE)

    # __bold__ -> **bold**
    out = re.sub(r"__([^_]+)__", r"**\1**", out)

    # Trim excessive blank lines around headings: at least one blank before # at line start (not first line)
    lines = out.split("\n")
    normalized: list[str] = []
    for i, line in enumerate(lines):
        if re.match(r"^#+\s", line) and normalized and normalized[-1].strip() != "":
            normalized.append("")
        normalized.append(line)
        if re.match(r"^#+\s", line):
            if i + 1 < len(lines) and lines[i + 1].strip() != "" and not re.match(
                r"^#+\s", lines[i + 1]
            ):
                pass  # ensure blank after header in second pass
    out = "\n".join(normalized)

    # One blank line after ATX headers
    parts = re.split(r"(^#+\s.*$)", out, flags=re.MULTILINE)
    rebuilt: list[str] = []
    for j, p in enumerate(parts):
        rebuilt.append(p)
        if re.match(r"^#+\s", p, flags=re.MULTILINE) and j + 1 < len(parts):
            nxt = parts[j + 1]
            if nxt and not nxt.startswith("\n"):
                rebuilt.append("\n")
    out = "".join(rebuilt)

    # Italic: prefer _word_ for simple *italic* not part of **
    def _italic_sub(m: re.Match[str]) -> str:
        inner = m.group(1)
        if "**" in inner:
            return m.group(0)
        return "_" + inner + "_"

    out = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", _italic_sub, out)

    return out.rstrip() + "\n"


def simple_concepts(text: str, n: int) -> list[str]:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_-]{5,}", text.lower())
    seen: set[str] = set()
    out: list[str] = []
    for w in words:
        if w in seen:
            continue
        seen.add(w)
        out.append(w)
        if len(out) >= n:
            break
    return out


# ---------------------------------------------------------------------------
# Chapter splitting (spindle prep)
# ---------------------------------------------------------------------------


def _strip_inline_flags(pattern: str) -> str:
    """Remove leading (?i) etc. so alternation compiles; use re.IGNORECASE instead."""
    p = pattern.strip()
    while True:
        if p.startswith("(?i)"):
            p = p[4:].lstrip()
            continue
        if p.startswith("(?-i)"):
            p = p[5:].lstrip()
            continue
        break
    return p


def combined_chapter_regex(patterns: list[str], fallback: str) -> str:
    patts = [_strip_inline_flags(p) for p in patterns if p and p.strip()]
    if not patts:
        patts = [_strip_inline_flags(fallback)]
    return "|".join(f"(?:{p})" for p in patts)


def split_into_chapters(full_text: str, combined_pattern: str) -> list[tuple[str, str]]:
    """Slice by chapter headers. Uses finditer so inner `(...)` in patterns cannot break splitting."""
    flags = re.MULTILINE | re.IGNORECASE
    try:
        rx = re.compile(combined_pattern, flags)
    except re.error as e:
        raise ValueError(f"Invalid chapter regex ({combined_pattern!r}): {e}") from e

    matches = list(rx.finditer(full_text))
    if not matches:
        return [("Full Document", full_text.strip())]

    chapters: list[tuple[str, str]] = []
    first_start = matches[0].start()
    intro = full_text[:first_start].strip()
    if intro:
        chapters.append(("Introduction", intro))

    for i, m in enumerate(matches):
        title = m.group(0).strip()
        body_start = m.end()
        body_end = matches[i + 1].start() if i + 1 < len(matches) else len(full_text)
        body = full_text[body_start:body_end].strip()
        chapters.append((title, body))
    return chapters


# ---------------------------------------------------------------------------
# PyMuPDF + OCR helpers (optional deps)
# ---------------------------------------------------------------------------


def pymupdf_extract_text(path: Path) -> tuple[str, int]:
    """Return (full_text, page_count). Empty if PyMuPDF unavailable or error."""
    if fitz is None:
        return "", 0
    try:
        doc = fitz.open(path_for_open(path))
    except Exception:
        return "", 0
    try:
        parts: list[str] = []
        for i in range(doc.page_count):
            parts.append(doc.load_page(i).get_text("text") or "")
        return "\n\n".join(parts).strip(), doc.page_count
    finally:
        doc.close()


def ocr_pil_image(img: Any, lang: str) -> str:
    if pytesseract is None:
        return ""
    try:
        return (pytesseract.image_to_string(img, lang=lang) or "").strip()
    except Exception:
        return ""


def ocr_image_path(path: Path, lang: str) -> tuple[str, list[str]]:
    warnings: list[str] = []
    if pytesseract is None or Image is None:
        warnings.append(
            "ocr_unavailable: pip install pytesseract pillow; install system tesseract-ocr (tesseract binary on PATH)"
        )
        return "", warnings
    try:
        with Image.open(path_for_open(path)) as im:
            rgb = im.convert("RGB")
            text = ocr_pil_image(rgb, lang)
        if not text:
            warnings.append("ocr_image: empty result (try other ocr_langs or image quality)")
        return text, warnings
    except Exception as e:
        warnings.append(f"ocr_image_failed: {e}")
        return "", warnings


def ocr_pdf_pages(path: Path, max_pages: int, dpi: int, lang: str) -> tuple[str, list[str]]:
    warnings: list[str] = []
    if fitz is None:
        warnings.append("ocr_pdf_unavailable: pip install pymupdf")
        return "", warnings
    if pytesseract is None or Image is None:
        warnings.append(
            "ocr_pdf_unavailable: pip install pytesseract pillow; install tesseract-ocr binary"
        )
        return "", warnings
    out_chunks: list[str] = []
    try:
        doc = fitz.open(path_for_open(path))
    except Exception as e:
        warnings.append(f"ocr_pdf_open_failed: {e}")
        return "", warnings
    try:
        n = min(doc.page_count, max(0, max_pages))
        if doc.page_count > max_pages:
            warnings.append(
                f"ocr_pdf: processing {n}/{doc.page_count} pages (ocr_max_pdf_pages cap)"
            )
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for i in range(n):
            try:
                page = doc.load_page(i)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                t = ocr_pil_image(img, lang)
                if t:
                    out_chunks.append(f"## Page {i + 1}\n\n{t}")
            except Exception as e:
                warnings.append(f"ocr_pdf_page_{i + 1}: {e}")
        return "\n\n".join(out_chunks).strip(), warnings
    finally:
        doc.close()


def length_balance_ratio(a: str, b: str) -> float | None:
    la, lb = len(a), len(b)
    if la < 20 or lb < 20:
        return None
    return min(la, lb) / max(la, lb)


# ---------------------------------------------------------------------------
# Engine
# ---------------------------------------------------------------------------

@dataclass
class ProcessResult:
    ok: bool
    source: str
    message: str = ""
    duration_s: float = 0.0
    outputs: list[str] = field(default_factory=list)
    sha256: str = ""
    warnings: list[str] = field(default_factory=list)


class ObsidianMasterEngine:
    def __init__(self, config: dict[str, Any]):
        self.config = config
        self._mid = None
        self._nlp = None
        if MarkItDown is not None:
            self._mid = MarkItDown()
        self._load_nlp()

    def _load_nlp(self) -> None:
        self._nlp = None
        if spacy is None:
            return
        model = self.config.get("nlp_settings", {}).get("model", "en_core_web_sm")
        try:
            self._nlp = spacy.load(model)
        except Exception:
            self._nlp = None

    def extract_tags(self, content_sample: str) -> set[str]:
        tags: set[str] = set(self.config.get("nlp_settings", {}).get("user_tags") or [])
        if self._nlp and content_sample:
            try:
                doc = self._nlp(content_sample[:50_000])
                want = set(self.config.get("nlp_settings", {}).get("extract_entities") or [])
                for ent in doc.ents:
                    if ent.label_ in want and ent.text.strip():
                        tags.add(ent.text.strip()[:100])
            except Exception:
                pass
        return tags

    def _plain_text_fallback(self, path: Path) -> tuple[str, list[str]]:
        warnings: list[str] = []
        raw = path.read_text(encoding="utf-8", errors="replace")
        if self._mid is None:
            warnings.append(
                "markitdown_unavailable: extracted as plain UTF-8 (install markitdown for PDF/DOCX)"
            )
        return raw.strip(), warnings

    def _docx_fallback(self, path: Path) -> tuple[str, list[str]]:
        warnings: list[str] = []
        if docx is None:
            raise RuntimeError("python-docx not installed; pip install python-docx (DOCX fallback)")
        try:
            d = docx.Document(path_for_open(path))
            paras = [(p.text or "").rstrip() for p in d.paragraphs]
            text = "\n\n".join([p for p in paras if p.strip()]).strip()
            if not text:
                warnings.append("docx_empty: no paragraph text extracted")
            return text, warnings
        except Exception as e:
            raise RuntimeError(f"docx_fallback_failed: {e}") from e

    def _json_fallback(self, path: Path) -> tuple[str, list[str]]:
        warnings: list[str] = []
        try:
            obj = json.loads(path.read_text(encoding="utf-8", errors="replace"))
            return "```json\n" + json.dumps(obj, indent=2, ensure_ascii=False) + "\n```\n", warnings
        except Exception as e:
            warnings.append(f"json_parse_failed: {e}")
            raw = path.read_text(encoding="utf-8", errors="replace")
            return "```text\n" + raw + "\n```\n", warnings

    def _csv_fallback(self, path: Path, max_rows: int = 200) -> tuple[str, list[str]]:
        warnings: list[str] = []
        try:
            with open(path_for_open(path), newline="", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows: list[list[str]] = []
                for i, row in enumerate(reader):
                    if i >= max_rows:
                        warnings.append(f"csv_truncated: max_rows={max_rows}")
                        break
                    rows.append([c.strip() for c in row])
        except Exception as e:
            raise RuntimeError(f"csv_read_failed: {e}") from e

        if not rows:
            warnings.append("csv_empty")
            return "", warnings

        # Normalize to rectangular
        width = max(len(r) for r in rows)
        for r in rows:
            if len(r) < width:
                r.extend([""] * (width - len(r)))

        header = rows[0]
        body = rows[1:] if len(rows) > 1 else []
        md = []
        md.append("| " + " | ".join(header) + " |")
        md.append("| " + " | ".join(["---"] * width) + " |")
        for r in body:
            md.append("| " + " | ".join(r) + " |")
        return "\n".join(md) + "\n", warnings

    def _xlsx_fallback(self, path: Path, max_rows: int = 200) -> tuple[str, list[str]]:
        warnings: list[str] = []
        if openpyxl is None:
            raise RuntimeError("openpyxl not installed; pip install openpyxl (XLSX fallback)")
        try:
            wb = openpyxl.load_workbook(path_for_open(path), data_only=True, read_only=True)
        except Exception as e:
            raise RuntimeError(f"xlsx_open_failed: {e}") from e

        try:
            out: list[str] = []
            for ws in wb.worksheets:
                out.append(f"## Sheet: {ws.title}")
                rows_out: list[list[str]] = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= max_rows:
                        warnings.append(f"xlsx_truncated:{ws.title}: max_rows={max_rows}")
                        break
                    rows_out.append([("" if v is None else str(v)) for v in row])
                if not rows_out:
                    out.append("_Empty sheet_\n")
                    continue
                width = max(len(r) for r in rows_out)
                for r in rows_out:
                    if len(r) < width:
                        r.extend([""] * (width - len(r)))
                header = rows_out[0]
                body = rows_out[1:] if len(rows_out) > 1 else []
                out.append("| " + " | ".join(header) + " |")
                out.append("| " + " | ".join(["---"] * width) + " |")
                for r in body:
                    out.append("| " + " | ".join(r) + " |")
                out.append("")
            return "\n".join(out).strip() + "\n", warnings
        finally:
            try:
                wb.close()
            except Exception:
                pass

    def _pptx_fallback(self, path: Path) -> tuple[str, list[str]]:
        warnings: list[str] = []
        if pptx is None:
            raise RuntimeError("python-pptx not installed; pip install python-pptx (PPTX fallback)")
        try:
            prs = pptx.Presentation(path_for_open(path))
        except Exception as e:
            raise RuntimeError(f"pptx_open_failed: {e}") from e

        lines: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"## Slide {i}")
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        slide_texts.append(t)
            if slide_texts:
                lines.extend(["- " + t.replace("\n", " ") for t in slide_texts])
            else:
                lines.append("_No text on slide_")
            lines.append("")
        out = "\n".join(lines).strip()
        if not out:
            warnings.append("pptx_empty")
        return out + "\n", warnings

    def _html_fallback(self, path: Path) -> tuple[str, list[str]]:
        warnings: list[str] = []
        raw = path.read_text(encoding="utf-8", errors="replace")
        # Extremely small, safe fallback: strip tags to text blocks.
        # This is not a full HTML→Markdown converter; it is a readability fallback.
        text = re.sub(r"(?is)<(script|style).*?>.*?</\\1>", "", raw)
        text = re.sub(r"(?s)<[^>]+>", " ", text)
        text = html.unescape(text)
        text = re.sub(r"[ \\t]+", " ", text)
        text = re.sub(r"\\n\\s*\\n\\s*\\n+", "\\n\\n", text)
        return text.strip(), warnings

    def _xml_fallback(self, path: Path) -> tuple[str, list[str]]:
        warnings: list[str] = []
        raw = path.read_text(encoding="utf-8", errors="replace")
        try:
            root = ElementTree.fromstring(raw)
        except Exception as e:
            warnings.append(f"xml_parse_failed: {e}")
            return "```xml\n" + raw + "\n```\n", warnings

        def walk(el: ElementTree.Element, depth: int = 0, out: list[str] | None = None) -> list[str]:
            if out is None:
                out = []
            indent = "  " * depth
            attrs = " ".join([f'{k}=\"{v}\"' for k, v in el.attrib.items()])
            line = f"{indent}- <{el.tag}" + (f" {attrs}" if attrs else "") + ">"
            out.append(line)
            if el.text and el.text.strip():
                out.append(f"{indent}  {el.text.strip()}")
            for ch in list(el):
                walk(ch, depth + 1, out)
            return out

        lines = walk(root)
        return "\n".join(lines).strip() + "\n", warnings

    def convert_document(self, path: Path) -> tuple[str, list[str], str]:
        """Return (markdown_body, warnings, content_source)."""
        pkm = self.config.get("pkm_settings", {})
        ss = self.config.get("system_settings", {})
        ds = self.config.get("designer_settings", {})
        warnings: list[str] = []
        suffix = path.suffix.lower()
        low = int(ss.get("low_content_threshold_chars", 200))
        ocr_on = bool(pkm.get("ocr_enabled", False))
        images_ocr = bool(pkm.get("images_ocr", True))
        ocr_lang = str(pkm.get("ocr_langs", "eng"))
        ocr_max_pages = int(pkm.get("ocr_max_pdf_pages", 40))
        ocr_dpi = int(pkm.get("ocr_dpi", 150))
        cross_threshold = float(pkm.get("ocr_comparison_threshold", 0.35))
        prefer_longer = bool(pkm.get("prefer_longer_text_on_divergence", True))
        pdf_pm = bool(pkm.get("pdf_pymupdf_text_fallback", True))
        cross = bool(pkm.get("pdf_cross_check_markitdown_vs_pymupdf", True))
        require_office = bool(pkm.get("require_markitdown_for_office_formats", False))
        fail_on_ocr_missing = bool(pkm.get("fail_if_ocr_enabled_but_unavailable", False))

        if ocr_on and fail_on_ocr_missing and not (pytesseract is not None and Image is not None and tesseract_is_available()):
            raise RuntimeError(
                "ocr_enabled=true but OCR dependencies are missing. "
                "Install Pillow + pytesseract + OS tesseract-ocr, or set fail_if_ocr_enabled_but_unavailable=false."
            )

        text = ""
        content_source = "unknown"

        if suffix in {".md", ".txt", ".markdown"}:
            text, w = self._plain_text_fallback(path)
            warnings.extend(w)
            content_source = "utf8_plain"

        elif suffix == ".json":
            text, w = self._json_fallback(path)
            warnings.extend(w)
            content_source = "stdlib_json"

        elif suffix == ".csv":
            text, w = self._csv_fallback(path)
            warnings.extend(w)
            content_source = "stdlib_csv"

        elif suffix == ".xlsx" and self._mid is None and not require_office:
            text, w = self._xlsx_fallback(path)
            warnings.extend(w)
            content_source = "openpyxl"

        elif suffix == ".pptx" and self._mid is None and not require_office:
            text, w = self._pptx_fallback(path)
            warnings.extend(w)
            content_source = "python-pptx"

        elif suffix in {".html", ".htm"}:
            text, w = self._html_fallback(path)
            warnings.extend(w)
            content_source = "stdlib_html_text"

        elif suffix == ".xml":
            text, w = self._xml_fallback(path)
            warnings.extend(w)
            content_source = "stdlib_xml"

        elif suffix == ".docx" and self._mid is None:
            # Fallback when MarkItDown is unavailable
            text, w = self._docx_fallback(path)
            warnings.extend(w)
            content_source = "python-docx"

        elif suffix in IMAGE_EXTENSIONS:
            content_source = "image_embed"
            if images_ocr and ocr_on:
                ocr_body, ow = ocr_image_path(path, ocr_lang)
                warnings.extend(ow)
                if ocr_body:
                    text = f"## Image OCR ({path.name})\n\n{ocr_body}"
                    content_source = "ocr_image"
                else:
                    text = (
                        f"![[{path.name}]]\n\n"
                        f"_OCR returned no text. Check tesseract install and `pkm_settings.ocr_langs`._\n"
                    )
                    warnings.append("image_ocr_empty")
            elif images_ocr and not ocr_on:
                text = f"![[{path.name}]]\n\n_Enable `pkm_settings.ocr_enabled` to OCR raster images._\n"
                warnings.append("image_ocr_skipped_disabled")
            else:
                text = f"![[{path.name}]]\n"

        elif suffix == ".pdf":
            md_text = ""
            if self._mid is not None:
                try:
                    result = self._mid.convert(str(path))
                    md_text = (result.text_content or "").strip()
                except Exception as e:
                    warnings.append(f"markitdown_pdf: {e}")
            pm_text = ""
            page_count = 0
            if pdf_pm:
                pm_text, page_count = pymupdf_extract_text(path)
                if not pm_text and fitz is None:
                    warnings.append("pymupdf_not_installed: pip install pymupdf for PDF text fallback")

            if cross and md_text and pm_text:
                lr = length_balance_ratio(md_text, pm_text)
                if lr is not None and lr < cross_threshold:
                    warnings.append(
                        f"cross_check: markitdown_vs_pymupdf length_ratio={lr:.3f} "
                        f"(threshold={cross_threshold}); sources diverge — review output"
                    )

            if prefer_longer:
                if len(md_text) >= len(pm_text):
                    text = md_text or pm_text
                    content_source = "markitdown" if md_text else "pymupdf"
                else:
                    text = pm_text
                    content_source = "pymupdf"
            else:
                text = md_text if md_text else pm_text
                content_source = "markitdown" if md_text else "pymupdf"

            if md_text and pm_text:
                content_source = "markitdown+pymupdf"

            # Image-heavy PDF: OCR fallback when text still short
            if len(text) < low and ocr_on:
                cap = ocr_max_pages if page_count == 0 else min(ocr_max_pages, page_count)
                ocr_body, ow = ocr_pdf_pages(path, cap, ocr_dpi, ocr_lang)
                warnings.extend(ow)
                if len(ocr_body) > len(text):
                    text = ocr_body
                    content_source = "ocr_pdf"
                    warnings.append("ocr_pdf_replaced_low_native_text")

            if not text.strip():
                raise RuntimeError(
                    "PDF produced no text (markitdown + pymupdf). "
                    "Enable pkm_settings.ocr_enabled and install pymupdf + pytesseract + tesseract binary, "
                    "or check the PDF."
                )

        else:
            if self._mid is None:
                if not require_office and suffix == ".docx":
                    text, w = self._docx_fallback(path)
                    warnings.extend(w)
                    content_source = "python-docx"
                elif not require_office and suffix == ".xlsx":
                    text, w = self._xlsx_fallback(path)
                    warnings.extend(w)
                    content_source = "openpyxl"
                elif not require_office and suffix == ".pptx":
                    text, w = self._pptx_fallback(path)
                    warnings.extend(w)
                    content_source = "python-pptx"
                else:
                    raise RuntimeError(
                        "markitdown not installed; pip install markitdown "
                        f"(required for {suffix})"
                    )
            try:
                if self._mid is not None:
                    result = self._mid.convert(str(path))
                    text = (result.text_content or "").strip()
            except Exception as e:
                raise RuntimeError(f"convert failed: {e}") from e
            if self._mid is not None:
                content_source = "markitdown"

        if len(text) < low:
            warnings.append(
                f"low_content: extracted text length {len(text)} < threshold {low} "
                "(short doc, image-heavy file, or OCR quality limit)"
            )

        text = apply_designer_linting(
            text,
            bool(ds.get("strict_markdown", True)),
            str(ds.get("list_marker", "-")),
        )
        return text, warnings, content_source

    def build_meta(
        self,
        source_name: str,
        file_hash: str,
        chapter_idx: int | None,
        chapter_title: str,
        tags: set[str],
        concepts: list[str],
        extra_warnings: list[str],
        content_source: str = "unknown",
    ) -> dict[str, Any]:
        template = dict(self.config.get("yaml_template") or {})
        template_tags = {t for t in (template.get("tags") or []) if isinstance(t, str)}
        merged_tags = sorted(template_tags | tags)
        meta: dict[str, Any] = {
            **template,
            "source": source_name,
            "sha256": file_hash,
            "chapter": chapter_idx,
            "chapter_title": chapter_title[:500],
            "tags": merged_tags,
            "concepts": concepts,
            "content_source": content_source[:120],
            "processed_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        }
        if extra_warnings:
            meta["warnings"] = extra_warnings
        _validate_meta_flat(meta)
        return meta

    def _note_basenames(
        self,
        stem_safe: str,
        chapters: list[tuple[str, str]],
        preserve_unicode: bool,
    ) -> list[str]:
        basenames: list[str] = []
        counts: dict[str, int] = {}
        for idx, (title, _body) in enumerate(chapters):
            slug = slugify_title(title, preserve_unicode=preserve_unicode)
            key = f"{stem_safe}_{idx:02d}_{slug}"
            if len(key) > 180:
                key = key[:180].rstrip("_")
            if key in counts:
                counts[key] += 1
                base = f"{key}_{counts[key]}"
            else:
                counts[key] = 0
                base = key
            basenames.append(base)
        return basenames

    def process_file(self, input_path: Path, output_dir: Path) -> ProcessResult:
        t0 = time.time()
        outputs: list[str] = []
        warnings: list[str] = []
        written_paths: list[Path] = []

        if not input_path.is_file():
            return ProcessResult(False, input_path.name, "not a file", time.time() - t0)

        h = file_sha256(input_path)
        ds = self.config.get("designer_settings", {})
        preserve = bool(ds.get("preserve_unicode_filenames", True))
        stem_safe = slugify_title(input_path.stem, max_len=60, preserve_unicode=preserve)

        try:
            content, conv_warnings, content_source = self.convert_document(input_path)
            warnings.extend(conv_warnings)
            for w in conv_warnings:
                logging.info("%s: %s", input_path.name, w)
        except Exception as e:
            return ProcessResult(False, input_path.name, str(e), time.time() - t0, sha256=h)

        pkm = self.config.get("pkm_settings", {})
        split = bool(pkm.get("split_chapters", True))
        designer_pat = str(ds.get("chapter_pattern", ""))
        patterns = list(pkm.get("chapter_patterns") or [])
        combined = combined_chapter_regex(patterns, designer_pat)

        if split:
            chapters = split_into_chapters(content, combined)
        else:
            chapters = [("Full Document", content)]

        basenames = self._note_basenames(stem_safe, chapters, preserve_unicode=preserve)
        index_basename = f"{stem_safe}_Index"

        n_concepts = int(
            self.config.get("nlp_settings", {}).get("extract_concepts_count", 5)
        )
        concepts = simple_concepts(content, n_concepts)
        tags = self.extract_tags(content)

        atomic = bool(self.config.get("system_settings", {}).get("atomic_writes", True))

        try:
            for idx, ((title, body), base) in enumerate(zip(chapters, basenames)):
                prev_l = f"[[{basenames[idx - 1]}]]" if idx > 0 else "START"
                next_l = f"[[{basenames[idx + 1]}]]" if idx + 1 < len(basenames) else "END"
                nav = (
                    f"\n\n---\n**Navigation**: {prev_l} | [[{index_basename}]] | {next_l}\n"
                )
                meta = self.build_meta(
                    input_path.name,
                    h,
                    idx + 1 if len(chapters) > 1 else None,
                    title,
                    tags,
                    concepts,
                    list(warnings) if idx == 0 else [],
                    content_source=content_source,
                )
                yaml_block = (
                    "---\n"
                    + yaml.safe_dump(
                        meta,
                        allow_unicode=True,
                        sort_keys=False,
                        default_flow_style=False,
                    )
                    + "---\n\n"
                )
                note_body = body + nav
                out_path = output_dir / f"{base}.md"
                full = yaml_block + note_body
                if atomic:
                    atomic_write_text(out_path, full)
                else:
                    out_path.parent.mkdir(parents=True, exist_ok=True)
                    out_path.write_text(full, encoding="utf-8")
                written_paths.append(out_path)
                outputs.append(str(out_path))

            index_lines = [f"# Index: {input_path.name}", ""]
            for b in basenames:
                index_lines.append(f"- [[{b}]]")
            index_meta = self.build_meta(
                input_path.name,
                h,
                None,
                "Index",
                tags,
                concepts,
                warnings,
                content_source=content_source,
            )
            index_meta["role"] = "index"
            yaml_block = (
                "---\n"
                + yaml.safe_dump(
                    index_meta,
                    allow_unicode=True,
                    sort_keys=False,
                    default_flow_style=False,
                )
                + "---\n\n"
            )
            index_path = output_dir / f"{index_basename}.md"
            index_body = "\n".join(index_lines) + "\n"
            if atomic:
                atomic_write_text(index_path, yaml_block + index_body)
            else:
                index_path.write_text(yaml_block + index_body, encoding="utf-8")
            written_paths.append(index_path)
            outputs.append(str(index_path))
        except Exception as e:
            for p in reversed(written_paths):
                try:
                    p.unlink()
                except OSError:
                    pass
            return ProcessResult(
                False, input_path.name, str(e), time.time() - t0, sha256=h
            )

        return ProcessResult(
            True,
            input_path.name,
            "",
            time.time() - t0,
            outputs,
            sha256=h,
            warnings=warnings,
        )


def _validate_meta_flat(meta: dict[str, Any]) -> None:
    """Reject nested structures that often break Obsidian property panels."""
    for k, v in meta.items():
        if isinstance(v, dict):
            raise ValueError(f"yaml field {k!r} must not be nested dict; flatten or stringify.")


# ---------------------------------------------------------------------------
# Multiprocessing (worker recycling)
# ---------------------------------------------------------------------------

_WORKER_ENGINE: ObsidianMasterEngine | None = None


def _pool_init(config_path: str) -> None:
    global _WORKER_ENGINE
    cfg = load_config(Path(config_path))
    _WORKER_ENGINE = ObsidianMasterEngine(cfg)


def _pool_process(task: tuple[str, str]) -> dict[str, Any]:
    inp, outd = task
    try:
        assert _WORKER_ENGINE is not None
        r = _WORKER_ENGINE.process_file(Path(inp), Path(outd))
        return {
            "ok": r.ok,
            "source": r.source,
            "message": r.message,
            "duration_s": r.duration_s,
            "outputs": r.outputs,
            "sha256": r.sha256,
            "warnings": r.warnings,
        }
    except Exception as e:
        logging.exception("Worker failed for %s", inp)
        h = ""
        try:
            h = file_sha256(Path(inp))
        except OSError:
            pass
        return {
            "ok": False,
            "source": Path(inp).name,
            "message": f"{type(e).__name__}: {e}",
            "duration_s": 0.0,
            "outputs": [],
            "sha256": h,
            "warnings": [],
        }


# ---------------------------------------------------------------------------
# Logs
# ---------------------------------------------------------------------------

def load_process_log(path: Path) -> dict[str, Any]:
    if not path.is_file():
        return {}
    with path.open(encoding="utf-8") as f:
        return json.load(f)


def save_process_log(path: Path, data: dict[str, Any]) -> None:
    with path.open("w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


ERROR_CSV_FIELDS = ("ts", "file", "error", "sha256")


def append_error_csv(path: Path, row: dict[str, str]) -> None:
    exists = path.is_file()
    normalized = {k: str(row.get(k, "")) for k in ERROR_CSV_FIELDS}
    with path.open("a", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=list(ERROR_CSV_FIELDS))
        if not exists:
            w.writeheader()
        w.writerow(normalized)


# ---------------------------------------------------------------------------
# Batch run
# ---------------------------------------------------------------------------

def iter_input_files(folder: Path, recursive: bool = False) -> list[Path]:
    exts = {
        ".pdf",
        ".docx",
        ".doc",
        ".pptx",
        ".xlsx",
        ".xls",
        ".html",
        ".htm",
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
    } | set(IMAGE_EXTENSIONS)
    out: list[Path] = []
    if recursive:
        for p in sorted(folder.rglob("*")):
            if p.is_file() and p.suffix.lower() in exts:
                out.append(p)
    else:
        for p in sorted(folder.iterdir()):
            if p.is_file() and p.suffix.lower() in exts:
                out.append(p)
    return out


def run_batch(
    config_path: Path,
    input_dir: Path,
    output_dir: Path,
    *,
    no_skip: bool = False,
    recursive: bool = False,
) -> int:
    cfg = load_config(config_path)
    ss = cfg["system_settings"]
    log_file = Path(ss.get("log_file", "process_log.json"))
    err_file = Path(ss.get("error_log", "error_log.csv"))
    cache = load_process_log(log_file)

    files = iter_input_files(input_dir, recursive=recursive)
    if not files:
        logging.warning("No supported files in %s", input_dir)
        return 1

    tier_name = ss.get("selected_tier", "safe")
    nproc = int(ss.get("core_tiers", {}).get(tier_name, 4))
    nproc = max(1, min(nproc, cpu_count() or 4))
    mode = ss.get("processing_mode", "parallel")
    max_tasks = int(ss.get("max_tasks_per_worker", 5))

    tasks: list[tuple[str, str]] = []
    for f in files:
        h = file_sha256(f)
        if not no_skip:
            prev = cache.get(h)
            if isinstance(prev, dict) and prev.get("status") == "success":
                logging.info("skip unchanged (sha256 hit): %s", f.name)
                continue
        tasks.append((str(f.resolve()), str(output_dir.resolve())))

    results: list[dict[str, Any]] = []
    if not tasks:
        logging.info("Nothing to do (all skipped).")
        return 0

    if mode == "parallel" and len(tasks) > 1:
        with Pool(
            processes=nproc,
            maxtasksperchild=max_tasks,
            initializer=_pool_init,
            initargs=(str(config_path.resolve()),),
        ) as pool:
            results = pool.map(_pool_process, tasks, chunksize=1)
    else:
        eng = ObsidianMasterEngine(cfg)
        for inp, outd in tasks:
            try:
                r = eng.process_file(Path(inp), Path(outd))
                results.append(
                    {
                        "ok": r.ok,
                        "source": r.source,
                        "message": r.message,
                        "duration_s": r.duration_s,
                        "outputs": r.outputs,
                        "sha256": r.sha256,
                        "warnings": r.warnings,
                    }
                )
            except Exception as e:
                logging.exception("Sequential worker failed for %s", inp)
                h = ""
                try:
                    h = file_sha256(Path(inp))
                except OSError:
                    pass
                results.append(
                    {
                        "ok": False,
                        "source": Path(inp).name,
                        "message": f"{type(e).__name__}: {e}",
                        "duration_s": 0.0,
                        "outputs": [],
                        "sha256": h,
                        "warnings": [],
                    }
                )

    for r in results:
        h = r.get("sha256", "")
        if r.get("ok"):
            cache[h] = {
                "status": "success",
                "file": r["source"],
                "time": r["duration_s"],
                "outputs": r.get("outputs", []),
                "warnings": r.get("warnings", []),
            }
        else:
            cache[h] = {
                "status": "failed",
                "file": r["source"],
                "error": r.get("message", ""),
            }
            append_error_csv(
                err_file,
                {
                    "ts": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
                    "file": r["source"],
                    "error": r.get("message", ""),
                    "sha256": h,
                },
            )

    save_process_log(log_file, cache)
    failed = sum(1 for r in results if not r.get("ok"))
    return 1 if failed else 0


# ---------------------------------------------------------------------------
# GUI (optional)
# ---------------------------------------------------------------------------


def run_gui(config_path: Path) -> None:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk

    cfg = load_config(config_path)
    eng = ObsidianMasterEngine(cfg)

    root = tk.Tk()
    root.title("Obsidian Designer Engine")

    src = tk.StringVar()
    dst = tk.StringVar()

    frm = ttk.Frame(root, padding=12)
    frm.pack(fill="both", expand=True)

    ttk.Label(frm, text="Source folder:").grid(row=0, column=0, sticky="w")
    ttk.Entry(frm, textvariable=src, width=48).grid(row=0, column=1)
    ttk.Button(
        frm,
        text="Browse",
        command=lambda: src.set(filedialog.askdirectory() or src.get()),
    ).grid(row=0, column=2)

    ttk.Label(frm, text="Output vault folder:").grid(row=1, column=0, sticky="w")
    ttk.Entry(frm, textvariable=dst, width=48).grid(row=1, column=1)
    ttk.Button(
        frm,
        text="Browse",
        command=lambda: dst.set(filedialog.askdirectory() or dst.get()),
    ).grid(row=1, column=2)

    pb = ttk.Progressbar(frm, length=420, mode="determinate")
    pb.grid(row=2, column=0, columnspan=3, pady=12)

    def on_run() -> None:
        s, d = src.get().strip(), dst.get().strip()
        if not s or not d:
            messagebox.showerror("Error", "Select source and output folders.")
            return
        files = iter_input_files(Path(s))
        pb["maximum"] = max(1, len(files))
        pb["value"] = 0
        for i, f in enumerate(files, start=1):
            r = eng.process_file(f, Path(d))
            if not r.ok:
                messagebox.showerror("Error", f"{r.source}: {r.message}")
                return
            pb["value"] = i
            root.update_idletasks()
        messagebox.showinfo("Done", "Batch complete.")

    ttk.Button(frm, text="Run", command=on_run).grid(row=3, column=1, pady=8)

    root.mainloop()


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    logging.basicConfig(
        level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s"
    )
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--config", type=Path, default=Path("config.json"))
    ap.add_argument("--input", type=Path, help="Input folder (documents)")
    ap.add_argument("--output", type=Path, help="Obsidian vault output folder")
    ap.add_argument("--gui", action="store_true", help="Tkinter UI")
    ap.add_argument("--doctor", action="store_true", help="Check dependencies and exit")
    ap.add_argument(
        "--no-skip",
        action="store_true",
        help="Reprocess all files even if process_log.json marks success for that hash",
    )
    ap.add_argument(
        "--recursive",
        action="store_true",
        help="Include supported files in subfolders (rglob)",
    )
    args = ap.parse_args()

    if args.doctor:
        cfg = load_config(args.config)
        code, rep = doctor_report(cfg)
        print(rep)
        raise SystemExit(code)

    if args.gui:
        run_gui(args.config)
        return

    if not args.input or not args.output:
        ap.error("--input and --output are required unless --gui")

    raise SystemExit(
        run_batch(
            args.config,
            args.input,
            args.output,
            no_skip=args.no_skip,
            recursive=args.recursive,
        )
    )


if __name__ == "__main__":
    main()
